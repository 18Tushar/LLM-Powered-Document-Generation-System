import chromadb
import os
import streamlit as st
import requests
import json
import random
import re
from langchain.vectorstores import Chroma
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.document_loaders import PyPDFLoader, TextLoader
from langchain.text_splitter import CharacterTextSplitter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO

# vLLM API endpoint
#VLLM_API = "http://localhost:8000/v1/completions"
VLLM_API_LAMA = "http://localhost:8000/v1/completions"

# Set page config for better appearance
st.set_page_config(
    layout="wide"
)

## RAG SETUP ##
# Initialize ChromaDB client and collection
CHROMA_DB_PATH = "./chroma_db"
chroma_client = chromadb.PersistentClient(path=CHROMA_DB_PATH)
chroma_collection = chroma_client.get_or_create_collection(
    name="ppt_knowledge")

# Load embedding model
embeddings = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2")

st.title("📚💻 Micropoint RAG-Powered AI Presentation Creator")

# Initialize session state variables
if "ppt_content" not in st.session_state:
    st.session_state.ppt_content = ""
if "topic" not in st.session_state:
    st.session_state.topic = ""
if "theme" not in st.session_state:
    st.session_state.theme = "minimal"
if "num_slides" not in st.session_state:
    st.session_state.num_slides = 8
if "content_per_slide" not in st.session_state:
    st.session_state.content_per_slide = 120
if "include_visuals" not in st.session_state:
    st.session_state.include_visuals = True

# File Upload Section
uploaded_file = st.file_uploader(
    "📂 Upload a document (PDF, TXT, DOCX)", type=[
        "pdf", "txt", "docx"])

if uploaded_file:
    with st.spinner("🔄 Processing document..."):
        # Save the uploaded file to a temporary path
        temp_file_path = os.path.join("/tmp", uploaded_file.name)

        with open(temp_file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        # Load and split text from the document
        if uploaded_file.name.endswith(".pdf"):
            loader = PyPDFLoader(temp_file_path)
        elif uploaded_file.name.endswith(".txt"):
            loader = TextLoader(temp_file_path)
        else:
            st.error("❌ Unsupported file format")
            st.stop()

        docs = loader.load()
        text_splitter = CharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        split_docs = text_splitter.split_documents(docs)

        # Store document embeddings in ChromaDB
        chroma_vectorstore = Chroma.from_documents(
            split_docs,
            embeddings,
            collection_name="ppt_knowledge",
            persist_directory=CHROMA_DB_PATH)

        st.success(
            f"✅ {len(split_docs)} document chunks stored in ChromaDB for retrieval!")

# Function to retrieve relevant data from ChromaDB


def retrieve_relevant_info(topic, num_results=5):  # Increased from 3 to 5
    """Retrieve relevant document snippets from ChromaDB."""
    if chroma_collection.count() == 0:
        return "No relevant documents found."

    results = chroma_collection.query(
        query_texts=[topic], n_results=num_results)

    if "documents" in results and results["documents"]:
        return "\n".join(results["documents"][0])
    else:
        return "No relevant content found."

# Function to generate presentation content using LLM + RAG


def generate_ppt_content(
        topic,
        num_slides,
        content_per_slide,
        include_visuals=True):
    """Generate PowerPoint content using LLM and ChromaDB knowledge retrieval."""

    # Retrieve knowledge from stored documents
    retrieved_info = retrieve_relevant_info(
        topic, num_results=5)  # Increased from 3 to 5

    # Updated prompt to emphasize detailed content and better formatting
    visual_instruction = ""
    if include_visuals:
        visual_instruction = """
        For each content slide, include ONE of the following visual element instructions (rotate between them for variety):

        CHART: [Specify chart type (bar, line, pie, etc.) and what data it should visualize]
        or
        IMAGE: [Describe an image that would enhance this slide]
        or
        DIAGRAM: [Suggest a diagram type (flow chart, mind map, etc.) and what it should show]

        Place this visual instruction on a new line after the bullet points using the format:
        VISUAL: [Your instruction]
        """

    # Enhanced prompt with stronger structure guidance, more detailed content,
    # and better formatting
    prompt = f"""
    Use the following reference information to generate accurate, detailed slides:
    {retrieved_info}

    Create a detailed, professional PowerPoint presentation on "{topic}" with exactly {num_slides} slides.

    Design a cohesive presentation that flows logically and engages the audience. Include:
    1. A compelling title slide with an impactful subtitle
    2. An agenda or overview slide that clearly outlines the presentation flow
    3. Detailed content slides with substantive information (5-6 bullet points per slide)
    4. At least one data or insights slide with specific facts/figures
    5. A comprehensive summary/conclusion slide with clear takeaways

    For EACH slide:
    - Create 5-6 detailed, informative bullet points per slide (more than in a typical presentation)
    - Each bullet should provide substantial information (30-40 words)
    - Use hierarchical bullet structure with main points and sub-points where appropriate
    - Include specific examples, case studies, or data points to reinforce each main idea
    - Total content per slide: approximately {content_per_slide} words
    - Ensure bullet points flow logically and build upon previous points
    - Use bullet points with • symbol for main points and - symbol for sub-points

    {visual_instruction}

    Use this EXACT FORMAT for your response:

    Slide 1: {topic}
    • [A compelling subtitle or tagline that captures the presentation's essence]
    • [Brief but comprehensive overview of the presentation's value proposition]
    • [Statement about the audience takeaways or benefits]

    Slide 2: [Detailed Agenda/Overview title]
    • [First key topic to be covered - be specific and include scope]
    • [Second key topic to be covered - be specific and include scope]
    • [Third key topic to be covered - be specific and include scope]
    • [Fourth key topic to be covered - be specific and include scope]
    • [Fifth key topic to be covered - be specific and include scope]
    • [Brief statement about the outcomes or takeaways expected from this presentation]
    VISUAL: [Visual instruction if applicable]

    Slide 3: [First main topic - be specific and engaging]
    • [Detailed main point 1 with substantial information - aim for 30-40 words with specific examples]
      - [Supporting detail or sub-point that elaborates on the main point]
      - [Additional supporting detail with specific facts or figures]
    • [Detailed main point 2 with substantial information - aim for 30-40 words with specific examples]
      - [Supporting detail or sub-point that elaborates on the main point]
      - [Additional supporting detail with specific facts or figures]
    • [Detailed main point 3 with substantial information - aim for 30-40 words with specific examples]
      - [Supporting detail or sub-point that elaborates on the main point]
    • [Detailed main point 4 with substantial information - aim for 30-40 words with specific examples]
    • [Detailed main point 5 with substantial information - aim for 30-40 words with specific examples]
    VISUAL: [Visual instruction if applicable]

    [Continue same pattern for remaining slides]

    Slide {num_slides}: [Comprehensive Conclusion/Summary/Next Steps title]
    • [Key takeaway 1 with actionable insight - be specific]
      - [Supporting point with specific recommendation]
    • [Key takeaway 2 with actionable insight - be specific]
      - [Supporting point with specific recommendation]
    • [Key takeaway 3 with actionable insight - be specific]
      - [Supporting point with specific recommendation]
    • [Key takeaway 4 with actionable insight - be specific]
    • [Key takeaway 5 with detailed call to action or next steps]
    • [Contact information or closing thought that reinforces main message]
    VISUAL: [Visual instruction if applicable]

    IMPORTANT: Do not include square brackets [] in your response. Replace them with actual content.
    Do not include quotation marks around any content.
    Ensure bullet points are substantive and detailed with facts, examples, and specific information.
    """

    # Set API Endpoint and Model Based on Selection
    if st.session_state.model_choice == "LLaMA-3-8B":
        model_name = "/data/Meta-Llama-3-8B-Instruct"
        api_url = VLLM_API_LAMA  
#    else:
#        model_name = "deepseek-ai/deepseek-r1-distill-llama-8b"  
#        api_url = VLLM_API

    payload = {
        "model": model_name,  # Dynamically select the model
        "prompt": prompt,
        "max_tokens": 2000,
        "temperature": 0.7
    }

    try:
        response = requests.post(VLLM_API_LAMA, headers={"Content-Type": "application/json"}, json=payload)

        if response.status_code == 200:
            try:
                raw_text = response.json()["choices"][0]["text"].strip()
                cleaned_text = re.sub(r'\[|\]', '', raw_text)
                cleaned_text = re.sub(r'["\'](.*?)["\']', r'\1', cleaned_text)
                return cleaned_text

            except (KeyError, IndexError):  # ✅ This must be inside the try block
                return "Error: Unexpected response format."

        else:
            return f"Error: {response.text}"

    except Exception as e:  # ✅ This must be aligned with `try:`
        return f"Error connecting to API: {str(e)}"

# Improved function to parse the generated content including nested bullet
# points
def parse_content(content):
    slides = []
    current_slide = None
    current_point = None

    # Split content by lines
    lines = content.split('\n')

    # Regex pattern to detect slide headers
    slide_pattern = re.compile(r'^slide\s+(\d+)[\s\:]+(.+)', re.IGNORECASE)

    for line_num, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue

        # Check if this is a slide title
        slide_match = slide_pattern.match(line)
        if slide_match:
            # Save previous slide if it exists
            if current_slide and current_slide.get("points"):
                slides.append(current_slide)

            # Extract title
            title = slide_match.group(2).strip()
            current_slide = {"title": title, "points": [], "visual": None}

        # Check for visual instructions
        elif current_slide and line.upper().startswith("VISUAL:"):
            visual_instruction = line[7:].strip()  # Remove "VISUAL:" prefix
            current_slide["visual"] = visual_instruction

        # Handle main bullet points
        elif current_slide and line.startswith("•"):
            # Remove the bullet point character
            point = line[1:].strip()
            point = re.sub(r'\[|\]|\"|\'', "", point)

            if point:  # Only add non-empty points
                point_obj = {"text": point, "subpoints": []}
                current_slide["points"].append(point_obj)
                current_point = point_obj

        # Handle sub-bullet points
        elif current_slide and current_point and line.startswith("-"):
            # Remove the bullet point character
            subpoint = line[1:].strip()
            subpoint = re.sub(r'\[|\]|\"|\'', "", subpoint)

            if subpoint:  # Only add non-empty subpoints
                current_point["subpoints"].append(subpoint)

        # Special case for numeric bullets like "1." or "1)"
        elif current_slide and re.match(r'^\d+[\.\)]', line):
            point = re.sub(r'^\d+[\.\)]', '', line).strip()
            # Clean the point
            point = re.sub(r'\[|\]|\"|\'', "", point)
            if point:
                point_obj = {"text": point, "subpoints": []}
                current_slide["points"].append(point_obj)
                current_point = point_obj

        # Alternative slide title detection (for cases without "Slide X:"
        # format)
        elif line_num < len(lines) - 1 and not any(line.startswith(symbol) for symbol in ["•", "-", "VISUAL:"]):
            next_line = lines[line_num + 1].strip()
            if (next_line.startswith("•") or re.match(
                    r'^\d+[\.\)]', next_line)) and len(line) < 60:
                # This looks like a slide title
                if current_slide and current_slide.get("points"):
                    slides.append(current_slide)
                current_slide = {"title": line, "points": [], "visual": None}

    # Add the last slide if it exists
    if current_slide and current_slide.get("points"):
        slides.append(current_slide)

    # If no slides were detected, try fallback parsing approach
    if not slides:
        return fallback_parse(content)

    return slides

# Fallback parsing logic for when the usual patterns aren't found

def fallback_parse(content):
    slides = []
    paragraphs = content.split('\n\n')

    # Process each paragraph to extract structured content
    for i, para in enumerate(paragraphs):
        para = para.strip()
        if not para:
            continue

        lines = para.split('\n')
        title_line = lines[0] if lines else "Slide Content"

        # Check if the paragraph looks like a slide with a title and content
        if len(lines) > 1:
            # First line is likely the title
            title = title_line

            # The rest are content or bullet points
            points = []
            visual = None

            current_point = None

            for line in lines[1:]:
                line = line.strip()
                if not line:
                    continue

                # Check for visual instruction
                if line.upper().startswith("VISUAL:"):
                    visual = line[7:].strip()

                # Handle main bullet points
                elif line.startswith("•"):
                    # Remove the bullet point character
                    point = line[1:].strip()
                    point = re.sub(r'\[|\]|\"|\'', "", point)

                    if point:  # Only add non-empty points
                        point_obj = {"text": point, "subpoints": []}
                        points.append(point_obj)
                        current_point = point_obj

                # Handle sub-bullet points
                elif current_point and line.startswith("-"):
                    # Remove the bullet point character
                    subpoint = line[1:].strip()
                    subpoint = re.sub(r'\[|\]|\"|\'', "", subpoint)

                    if subpoint:  # Only add non-empty subpoints
                        current_point["subpoints"].append(subpoint)

                # Handle numeric bullets
                elif re.match(r'^\d+[\.\)]', line):
                    point = re.sub(r'^\d+[\.\)]', '', line).strip()
                    point = re.sub(r'\[|\]|\"|\'', "", point)

                    if point:
                        point_obj = {"text": point, "subpoints": []}
                        points.append(point_obj)
                        current_point = point_obj

                # Plain text (probably part of the previous point)
                elif current_point:
                    line = re.sub(r'\[|\]|\"|\'', "", line)
                    if line:
                        # Append to the current point text
                        current_point["text"] += " " + line

            if points:  # Only add slides with content
                slides.append({
                    "title": title,
                    "points": points,
                    "visual": visual
                })

    # If still no slides, create a minimal default set
    if not slides:
        slides = [{"title": "Generated Presentation",
                   "points": [{"text": "Your presentation content will appear here.",
                               "subpoints": []},
                              {"text": "Each slide should have a title followed by bullet points.",
                               "subpoints": []},
                              {"text": "Try adjusting your prompt for better results.",
                               "subpoints": []}],
                   "visual": None}]

    return slides


# Define modern theme color sets
THEMES = {
    "minimal": {
        "primary": RGBColor(30, 30, 30),  # Near black
        "secondary": RGBColor(248, 249, 250),  # Off-white
        "accent1": RGBColor(79, 134, 247),  # Bright blue
        "accent2": RGBColor(210, 210, 210),  # Light gray
        "text": RGBColor(40, 40, 40)  # Dark gray
    },
    "gradient": {
        "primary": RGBColor(67, 67, 186),  # Deep blue
        "secondary": RGBColor(252, 252, 255),  # White blue tint
        "accent1": RGBColor(141, 85, 232),  # Purple
        "accent2": RGBColor(95, 209, 249),  # Light blue
        "text": RGBColor(33, 33, 33)  # Dark gray
    },
    "vibrant": {
        "primary": RGBColor(13, 71, 161),  # Deep blue
        "secondary": RGBColor(250, 250, 250),  # White
        "accent1": RGBColor(0, 176, 155),  # Teal
        "accent2": RGBColor(255, 82, 82),  # Coral red
        "text": RGBColor(40, 40, 40)  # Dark gray
    },
    "dark": {
        "primary": RGBColor(30, 30, 35),  # Dark gray
        "secondary": RGBColor(18, 18, 23),  # Nearly black
        "accent1": RGBColor(255, 122, 69),  # Orange
        "accent2": RGBColor(133, 216, 206),  # Mint
        "text": RGBColor(240, 240, 240)  # Off-white
    }
}


def create_modern_ppt(slides_data, title, theme="minimal"):
    try:
        prs = Presentation()

        # Get theme colors
        theme_colors = THEMES.get(theme, THEMES["minimal"])

        # Set slide dimensions for 16:9 aspect ratio
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Function to add a modern background to slides based on theme
        def add_modern_background(slide, theme_colors, slide_type="content"):
            # Add background shape for the entire slide
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0,
                prs.slide_width,
                prs.slide_height
            )
            background.fill.solid()
            background.fill.fore_color.rgb = theme_colors["secondary"]
            background.line.fill.background()

            if theme == "gradient":
                # Add gradient accent in top-right corner
                accent_corner = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    prs.slide_width - Inches(5),
                    -Inches(2),
                    Inches(6),
                    Inches(6)
                )
                accent_corner.fill.solid()
                accent_corner.fill.fore_color.rgb = theme_colors["accent2"]
                accent_corner.line.fill.background()

                # Set transparency
                accent_corner.fill.transparency = 0.7

            elif theme == "minimal":
                # Add thin accent bar on bottom
                accent = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    0, prs.slide_height - Inches(0.2),
                    prs.slide_width,
                    Inches(0.2)
                )
                accent.fill.solid()
                accent.fill.fore_color.rgb = theme_colors["accent1"]
                accent.line.fill.background()

            elif theme == "vibrant":
                # Add diagonal accent element
                accent = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_TRIANGLE,
                    0, 0,
                    Inches(2),
                    Inches(7.5)
                )
                accent.fill.solid()
                accent.fill.fore_color.rgb = theme_colors["accent1"]
                accent.line.fill.background()

                small_accent = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(2), 0,
                    Inches(11.33),
                    Inches(0.3)
                )
                small_accent.fill.solid()
                small_accent.fill.fore_color.rgb = theme_colors["accent2"]
                small_accent.line.fill.background()

            elif theme == "dark":
                # Add subtle design element
                accent_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    -Inches(2), Inches(3),
                    Inches(15),
                    Inches(2)
                )
                accent_shape.fill.solid()
                accent_shape.fill.fore_color.rgb = theme_colors["primary"]
                accent_shape.line.fill.background()

                # Add a subtle light bar
            #    light_bar = slide.shapes.add_shape(
            #        MSO_SHAPE.RECTANGLE,
            #        0, Inches(1),
            #        prs.slide_width,
            #        Inches(0.05)
            #    )
            #    light_bar.fill.solid()
            #    light_bar.fill.fore_color.rgb = theme_colors["accent1"]
            #    light_bar.line.fill.background()

            # Special title slide treatment
            if slide_type == "title":
                if theme == "minimal":
                    # Bold colored shape for title slide
                    title_accent = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        0, 0,
                        prs.slide_width,
                        Inches(3.5)
                    )
                    title_accent.fill.solid()
                    title_accent.fill.fore_color.rgb = theme_colors["primary"]
                    title_accent.line.fill.background()

                elif theme == "gradient":
                    # Full gradient overlay for title
                    title_accent = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        0, 0,
                        prs.slide_width,
                        prs.slide_height
                    )
                    title_accent.fill.solid()
                    title_accent.fill.fore_color.rgb = theme_colors["primary"]
                    title_accent.line.fill.background()
                    title_accent.fill.transparency = 0.85

                    # Add design element
                    design_el = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(6), Inches(2),
                        Inches(10),
                        Inches(10)
                    )
                    design_el.fill.solid()
                    design_el.fill.fore_color.rgb = theme_colors["accent1"]
                    design_el.line.fill.background()
                    design_el.fill.transparency = 0.8

                elif theme == "vibrant":
                    # Split design for title
                    title_accent = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        0, 0,
                        Inches(5),
                        prs.slide_height
                    )
                    title_accent.fill.solid()
                    title_accent.fill.fore_color.rgb = theme_colors["primary"]
                    title_accent.line.fill.background()

                elif theme == "dark":
                    # Dramatic dark title slide
                    title_accent = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        0, 0,
                        prs.slide_width,
                        prs.slide_height
                    )
                    title_accent.fill.solid()
                    title_accent.fill.fore_color.rgb = theme_colors["primary"]
                    title_accent.line.fill.background()

                    # Add design element
                    light_bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(4), Inches(3.5),
                        Inches(5),
                        Inches(0.1)
                    )
                    light_bar.fill.solid()
                    light_bar.fill.fore_color.rgb = theme_colors["accent1"]
                    light_bar.line.fill.background()

            # Return content margins
            left_margin = Inches(0.6)
            top_margin = Inches(0.6)
            content_width = prs.slide_width - Inches(2.0)
            content_height = prs.slide_height - Inches(2.0)

            return left_margin, top_margin, content_width, content_height

        # Title Slide
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Add custom background for title slide
        left_margin, top_margin, content_width, content_height = add_modern_background(
            slide, theme_colors, "title")

        # Add title with modern styling
        title_top = Inches(1.5) if theme != "vibrant" else Inches(2)
        title_left = Inches(1) if theme != "vibrant" else Inches(5.5)
        title_width = content_width if theme != "vibrant" else Inches(7)

        title_shape = slide.shapes.add_textbox(
            title_left, title_top,
            title_width, Inches(1.8)
        )
        title_frame = title_shape.text_frame
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        p = title_frame.add_paragraph()

        # Clean title of any brackets or quotes that might be left over
        clean_title = re.sub(r'\[|\]|\"|\'', "", title).upper()
        p.text = clean_title

        p.alignment = PP_ALIGN.LEFT if theme != "minimal" else PP_ALIGN.CENTER
        p.font.size = Pt(54)
        p.font.bold = True

        # Set title text color based on theme and position
        if theme == "dark" or (theme == "minimal" and title_top < Inches(3.5)):
            p.font.color.rgb = RGBColor(
                255, 255, 255)  # White for dark backgrounds
        else:
            p.font.color.rgb = theme_colors["primary"]

        # Add subtitle with modern positioning
        subtitle_top = title_top + Inches(1.8)
        subtitle_shape = slide.shapes.add_textbox(
            title_left, subtitle_top,
            title_width, Inches(1)
        )
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.word_wrap = True

        p = subtitle_frame.add_paragraph()
        if slides_data and slides_data[0]["points"]:
            subtitle_text = slides_data[0]["points"][0]["text"] if isinstance(
                slides_data[0]["points"][0], dict) else slides_data[0]["points"][0]
            # Clean subtitle of any brackets or quotes
            subtitle_text = re.sub(r'\[|\]|\"|\'', "", subtitle_text)
            p.text = subtitle_text
        else:
            p.text = "A modern presentation generated with AI"
        p.alignment = PP_ALIGN.LEFT if theme != "minimal" else PP_ALIGN.CENTER
        p.font.size = Pt(28)

        # Set subtitle text color based on theme and position
        if theme == "dark":
            # Light gray for dark backgrounds
            p.font.color.rgb = RGBColor(230, 230, 230)
        elif theme == "minimal":
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black for minimal theme
        else:
            p.font.color.rgb = theme_colors["accent1"]

       # if theme == "dark" or (theme == "minimal" and subtitle_top < Inches(3.5)):
        #    p.font.color.rgb = RGBColor(230, 230, 230)  # Light gray for dark backgrounds
        # else:
         #   p.font.color.rgb = theme_colors["accent1"]

        # Remove the title slide data from the slides_data if it exists
        if slides_data and "title" in slides_data[0]["title"].lower():
            slides_data = slides_data[1:]

        # Content Slides
        for i, slide_data in enumerate(slides_data):
            slide = prs.slides.add_slide(
                prs.slide_layouts[1])  # Content layout

            # Add modern background
            left_margin, top_margin, content_width, content_height = add_modern_background(
                slide, theme_colors)

            # Add slide number with modern styling
            slide_number = slide.shapes.add_textbox(
                prs.slide_width - Inches(1.5),
                prs.slide_height - Inches(0.7),
                Inches(1),
                Inches(0.4)
            )
            slide_number_frame = slide_number.text_frame
            p = slide_number_frame.add_paragraph()
            p.text = f"{i+1}/{len(slides_data)}"
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(14)
            p.font.color.rgb = theme_colors["accent2"]

            # Add title with better positioning and styling
            title_shape = slide.shapes.add_textbox(
                left_margin, top_margin,
                Inches(10), Inches(0.8)
            )
            title_frame = title_shape.text_frame
            title_frame.word_wrap = True

            p = title_frame.add_paragraph()
            title_text = slide_data["title"]
            # Clean up title text from formatting remnants
            if ":" in title_text and title_text.split(
                    ":", 1)[0].strip().lower().startswith("slide"):
                title_text = title_text.split(":", 1)[1].strip()
            # Clean title of any brackets or quotes
            title_text = re.sub(r'\[|\]|\"|\'', "", title_text)
            p.text = title_text
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = theme_colors["primary"]

            # Determine if we have a visual instruction and position content
            # accordingly
            has_visual = slide_data.get("visual") is not None

            # Define visual dimensions and position
            visual_width = Inches(3)
            visual_height = Inches(3)
            visual_left = content_width - \
                visual_width + left_margin - Inches(0)
            visual_top = top_margin + Inches(1.5)

            # Adjust content width and position based on visual presence
            if has_visual:
                # Content goes on left side with reduced width to accommodate
                # visual on right
                content_left = left_margin
                content_top = title_top + Inches(0.0)
                content_box_width = visual_left - left_margin - \
                    Inches(0)  # Space between content and visual
            else:
                # Content takes full width when no visual is present
                content_left = left_margin
                content_top = top_margin + Inches(1)
                content_box_width = content_width - Inches(0.5)

            # Content area with adjusted width and position
            content_shape = slide.shapes.add_textbox(
                content_left, content_top,
                content_box_width, content_height - Inches(1.5)
            )
            content_frame = content_shape.text_frame
            content_frame.word_wrap = True

            # Make bullet points more visual
            for point in slide_data["points"]:
                if isinstance(point, dict):
                    # Main point
                    p = content_frame.add_paragraph()
                    p.text = "❖ " + point["text"]
                    p.level = 0
                    # Slightly smaller font for more space
                    p.font.size = Pt(22)
                    p.font.color.rgb = theme_colors["text"]
                    p.space_after = Pt(8)  # Add spacing between points

                    # Add subpoints if they exist
                    for subpoint in point.get("subpoints", []):
                        p = content_frame.add_paragraph()
                        p.text = "  ‣ " + subpoint
                        p.level = 1
                        p.font.size = Pt(18)  # Smaller font for subpoints
                        p.font.color.rgb = theme_colors["text"]
                        p.space_after = Pt(4)  # Less spacing for subpoints
                else:
                    # Handle case where point is a simple string
                    p = content_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = Pt(22)
                    p.font.color.rgb = theme_colors["text"]
                    p.space_after = Pt(8)

            # Add visual placeholder based on visual instruction
            if has_visual:
                visual_desc = slide_data["visual"]

                visual_type = "CHART"
                if "IMAGE:" in visual_desc.upper():
                    visual_type = "IMAGE"
                elif "DIAGRAM:" in visual_desc.upper():
                    visual_type = "DIAGRAM"

                # Add visual placeholder shape
                if visual_type == "CHART":
                    chart_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        visual_left, visual_top,
                        visual_width, visual_height
                    )
                    chart_shape.fill.solid()
                    chart_shape.fill.fore_color.rgb = theme_colors["accent2"]
                    chart_shape.line.color.rgb = theme_colors["accent1"]

                    # Add icon text
                    chart_text = slide.shapes.add_textbox(
                        visual_left + Inches(1.5), visual_top + Inches(1),
                        Inches(1.5), Inches(0.5)
                    )
                    p = chart_text.text_frame.add_paragraph()
                    p.text = "📊"
                    p.font.size = Pt(32)
                    p.alignment = PP_ALIGN.CENTER

                elif visual_type == "IMAGE":
                    image_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        visual_left, visual_top,
                        visual_width, visual_height
                    )
                    image_shape.fill.solid()
                    image_shape.fill.fore_color.rgb = theme_colors["accent2"]
                    image_shape.line.color.rgb = theme_colors["accent1"]

                    # Add icon text
                    image_text = slide.shapes.add_textbox(
                        visual_left + Inches(1.5), visual_top + Inches(1),
                        Inches(1.5), Inches(0.5)
                    )
                    p = image_text.text_frame.add_paragraph()
                  #  p.text = ""
                    p.font.size = Pt(32)
                    p.alignment = PP_ALIGN.CENTER

                elif visual_type == "DIAGRAM":
                    diagram_shape = slide.shapes.add_shape(
                        MSO_SHAPE.FLOWCHART_PROCESS,
                        visual_left, visual_top,
                        visual_width, visual_height
                    )
                    diagram_shape.fill.solid()
                    diagram_shape.fill.fore_color.rgb = theme_colors["accent2"]
                    diagram_shape.line.color.rgb = theme_colors["accent1"]

                    # Add icon text
                    diagram_text = slide.shapes.add_textbox(
                        visual_left + Inches(1.5), visual_top + Inches(1),
                        Inches(1.5), Inches(0.5)
                    )
                    p = diagram_text.text_frame.add_paragraph()
                    p.text = "📈"
                    p.font.size = Pt(32)
                    p.alignment = PP_ALIGN.CENTER

                # Add visual description below the shape
                desc_text = slide.shapes.add_textbox(
                    visual_left, visual_top + visual_height + Inches(0.1),
                    visual_width, Inches(0.6)
                )
                desc_text.text_frame.word_wrap = True
                p = desc_text.text_frame.add_paragraph()
                # Remove prefix like "CHART:"
                p.text = re.sub(r'^\w+:\s*', '', visual_desc)
                p.font.size = Pt(12)
                p.font.italic = True
                p.font.color.rgb = theme_colors["text"]
                p.alignment = PP_ALIGN.CENTER

        # Save to BytesIO for download
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io

    except Exception as e:
        st.error(f"Error creating PowerPoint: {str(e)}")
        return None


# App Layout and UI

# Model Selection Dropdown
st.session_state.model_choice = st.selectbox(
    "Choose Model:",
    options=["LLaMA-3-8B", "DeepSeek-R1-Distill-Llama-8B"],
    index=0  # Default selection
)

with st.container():
    st.markdown("#### 📝Enter Your Presentation Title")

    # User inputs
    col1, col2 = st.columns(2)

    with col1:
        st.session_state.topic = st.text_area(
            "Presentation Topic",
            value=st.session_state.topic,
            height=100,
            placeholder="Enter your presentation topic here...")

        st.session_state.num_slides = st.slider(
            "Number of Slides",
            min_value=5,
            max_value=15,
            value=st.session_state.num_slides)

        # Theme selection
        st.session_state.theme = st.selectbox(
            "Presentation Theme", options=list(
                THEMES.keys()), format_func=lambda x: x.capitalize(), index=list(
                THEMES.keys()).index(
                st.session_state.theme))

    with col2:
        st.session_state.content_per_slide = st.slider(
            "Content Per Slide (words)",
            min_value=80,
            max_value=200,
            value=st.session_state.content_per_slide)

        st.session_state.include_visuals = st.checkbox(
            "Include Visual Placeholders",
            value=st.session_state.include_visuals)

        # Generate button
        generate_btn = st.button(
            " Generate Presentation",
            use_container_width=True,
            type="primary")

# Generate presentation content when button is clicked
if generate_btn:
    with st.spinner("🤖 Generating presentation content..."):
        st.session_state.ppt_content = generate_ppt_content(
            st.session_state.topic,
            st.session_state.num_slides,
            st.session_state.content_per_slide,
            st.session_state.include_visuals
        )

    if "Error" in st.session_state.ppt_content:
        st.error(st.session_state.ppt_content)
    else:
        st.success("✅ Presentation content generated!")

# Display and download if content exists
if st.session_state.ppt_content and "Error" not in st.session_state.ppt_content:
    with st.expander("📑 View Presentation Content", expanded=True):
        st.markdown(st.session_state.ppt_content)

    # Parse content and create PPT
    slides_data = parse_content(st.session_state.ppt_content)

    with st.spinner("🔨 Building PowerPoint file..."):
        ppt_file = create_modern_ppt(
            slides_data,
            st.session_state.topic,
            st.session_state.theme)

    if ppt_file:
        # Add download button for the generated PPT
        st.download_button(
            label="📥 Download PowerPoint",
            data=ppt_file,
            file_name=f"{st.session_state.topic.replace(' ', '_')[:30]}_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True)

# Add footer information
st.markdown("---")
# st.markdown("")
# st.write("Retrieved Info from ChromaDB:", retrieved_info)
